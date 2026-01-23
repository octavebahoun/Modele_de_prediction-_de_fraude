from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()

    # Layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]

    # --- SLIDE 1: TITLE ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Détection de Fraude Financière"
    subtitle.text = "Analyse et Modélisation avec Random Forest\nJanvier 2026"

    # --- SLIDE 2: INTRODUCTION ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Introduction et Objectifs"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Objectif principal : Identifier automatiquement les transactions frauduleuses."
    p = tf.add_paragraph()
    p.text = "Points clés :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sécurité des transactions"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Réduction des pertes financières"
    p.level = 2

    # --- SLIDE 3: DATASET ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Exploration du Dataset"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Description des données :"
    p = tf.add_paragraph()
    p.text = "2000 transactions enregistrées"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Variables : montant (amount), heure (hour), transactions du jour, changement de lieu/appareil."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Cible (Label) : is_fraud (0 pour normal, 1 pour fraude)."
    p.level = 1

    # --- SLIDE 4: CORRELATION ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Analyse des Corrélations"
    if os.path.exists("outputs/correlation_matrix.png"):
        slide.shapes.add_picture("outputs/correlation_matrix.png", Inches(1), Inches(1.5), height=Inches(5))

    # --- SLIDE 5: METHODOLOGY ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Méthodologie"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Modèle utilisé : Random Forest Classifier"
    p = tf.add_paragraph()
    p.text = "Configuration :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "100 arbres de décision"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Répartition 80% entraînement / 20% test"
    p.level = 2

    # --- SLIDE 6: FEATURE IMPORTANCE ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Importance des Variables"
    if os.path.exists("outputs/feature_importance.png"):
        slide.shapes.add_picture("outputs/feature_importance.png", Inches(1), Inches(1.5), height=Inches(5))

    # --- SLIDE 7: PERFORMANCE ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Résultats et Matrice de Confusion"
    if os.path.exists("outputs/confusion_matrix.png"):
        slide.shapes.add_picture("outputs/confusion_matrix.png", Inches(1), Inches(1.5), height=Inches(5))

    # --- SLIDE 8: CONCLUSION ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion et Perspectives"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Le modèle présente une excellente capacité de détection."
    p = tf.add_paragraph()
    p.text = "Prochaines étapes :"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Tests sur des données en temps réel"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Optimisation des hyperparamètres"
    p.level = 2

    # Save
    if not os.path.exists("presentation"):
        os.makedirs("presentation")
    
    save_path = "presentation/presentation_fraude.pptx"
    prs.save(save_path)
    print(f"Presentation saved to {save_path}")

if __name__ == "__main__":
    create_presentation()
